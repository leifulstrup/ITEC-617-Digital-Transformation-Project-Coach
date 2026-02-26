# /// script
# requires-python = ">=3.12"
# dependencies = ["python-pptx>=1.0.2"]
# ///
"""
Extract content from a PowerPoint presentation for AI-powered feedback.

Reads .pptx files from student-workspace/slides/ and produces:
  - student-workspace/extracted/presentation-content.md  (structured markdown)
  - student-workspace/extracted/slide-images/             (any embedded images)

Usage:
    uv run extract_presentation.py
    uv run extract_presentation.py path/to/specific.pptx
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


SLIDES_DIR = Path("student-workspace/slides")
OUTPUT_DIR = Path("student-workspace/extracted")
IMAGES_DIR = OUTPUT_DIR / "slide-images"
OUTPUT_MD = OUTPUT_DIR / "presentation-content.md"


def find_pptx_file(explicit_path: str | None = None) -> Path:
    """Find the PowerPoint file to extract."""
    if explicit_path:
        path = Path(explicit_path)
        if not path.exists():
            print(f"Error: File not found: {path}")
            sys.exit(1)
        if path.suffix.lower() != ".pptx":
            print(f"Error: Expected .pptx file, got: {path.suffix}")
            sys.exit(1)
        return path

    pptx_files = list(SLIDES_DIR.glob("*.pptx"))
    if not pptx_files:
        print(f"Error: No .pptx files found in {SLIDES_DIR}/")
        print("Place your PowerPoint file in the student-workspace/slides/ directory.")
        sys.exit(1)
    if len(pptx_files) > 1:
        print(f"Found multiple .pptx files in {SLIDES_DIR}/:")
        for f in pptx_files:
            print(f"  - {f.name}")
        print("Using the most recently modified file.")
        pptx_files.sort(key=lambda p: p.stat().st_mtime, reverse=True)

    return pptx_files[0]


def extract_table_text(table) -> str:
    """Extract text from a table shape as a markdown table."""
    rows = []
    for row_idx, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
        if row_idx == 0:
            rows.append("| " + " | ".join(["---"] * len(row.cells)) + " |")
    return "\n".join(rows)


def extract_shape_text(shape) -> str:
    """Extract text content from a shape, handling different shape types."""
    parts = []

    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        parts.append(extract_table_text(shape.table))
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            child_text = extract_shape_text(child_shape)
            if child_text:
                parts.append(child_text)
    elif hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                # Detect bullet level for indentation
                level = paragraph.level or 0
                indent = "  " * level
                if level > 0:
                    parts.append(f"{indent}- {text}")
                else:
                    parts.append(text)
    elif hasattr(shape, "text") and shape.text.strip():
        parts.append(shape.text.strip())

    return "\n".join(parts)


def extract_images(slide, slide_num: int) -> list[str]:
    """Extract embedded images from a slide, save to disk, return filenames."""
    IMAGES_DIR.mkdir(parents=True, exist_ok=True)
    saved = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            filename = f"slide-{slide_num:02d}-{shape.shape_id}.{ext}"
            filepath = IMAGES_DIR / filename
            filepath.write_bytes(image.blob)
            saved.append(filename)

    return saved


def extract_presentation(pptx_path: Path) -> str:
    """Extract full presentation content as structured markdown."""
    prs = Presentation(str(pptx_path))

    lines = [
        "# Presentation Content",
        "",
        f"**Source**: `{pptx_path.name}`",
        f"**Slides**: {len(prs.slides)}",
        "",
        "---",
        "",
    ]

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    if slide_width and slide_height:
        width_in = slide_width / 914400
        height_in = slide_height / 914400
        lines.append(f"**Slide dimensions**: {width_in:.1f}\" x {height_in:.1f}\"")
        lines.append("")

    total_images = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {slide_num}")
        lines.append("")

        # Extract slide layout name if available
        if slide.slide_layout and slide.slide_layout.name:
            lines.append(f"*Layout: {slide.slide_layout.name}*")
            lines.append("")

        # Separate title from body content
        title_text = None
        body_parts = []

        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                # This is the title placeholder
                if hasattr(shape, "text") and shape.text.strip():
                    title_text = shape.text.strip()
            else:
                text = extract_shape_text(shape)
                if text:
                    body_parts.append(text)

        if title_text:
            lines.append(f"### {title_text}")
            lines.append("")

        if body_parts:
            for part in body_parts:
                lines.append(part)
                lines.append("")

        # Extract embedded images
        images = extract_images(slide, slide_num)
        if images:
            total_images += len(images)
            lines.append("**Embedded images:**")
            for img in images:
                lines.append(f"- `slide-images/{img}`")
            lines.append("")

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append("**Speaker Notes:**")
                lines.append("")
                lines.append(f"> {notes_text}")
                lines.append("")

        lines.append("---")
        lines.append("")

    # Summary
    lines.append("## Extraction Summary")
    lines.append("")
    lines.append(f"- **Total slides**: {len(prs.slides)}")
    lines.append(f"- **Embedded images extracted**: {total_images}")
    lines.append(f"- **Source file**: `{pptx_path.name}`")

    return "\n".join(lines)


def main():
    explicit_path = sys.argv[1] if len(sys.argv) > 1 else None
    pptx_path = find_pptx_file(explicit_path)

    print(f"Extracting content from: {pptx_path.name}")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    content = extract_presentation(pptx_path)
    OUTPUT_MD.write_text(content, encoding="utf-8")

    print(f"  Presentation content:  {OUTPUT_MD}")

    if IMAGES_DIR.exists() and any(IMAGES_DIR.iterdir()):
        image_count = len(list(IMAGES_DIR.iterdir()))
        print(f"  Embedded images ({image_count}): {IMAGES_DIR}/")

    print()
    print("Next steps:")
    print("  1. Review the extracted content in student-workspace/extracted/")
    print("  2. For visual slide analysis, also export your presentation as PDF:")
    print("     PowerPoint > File > Export > Create PDF")
    print("     Save the PDF to student-workspace/slides/")
    print("  3. Consult executive personas or run /evaluate for feedback")


if __name__ == "__main__":
    main()
